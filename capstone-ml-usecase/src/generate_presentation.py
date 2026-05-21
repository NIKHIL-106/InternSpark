from pptx import Presentation
from pptx.util import Inches, Pt
from fpdf import FPDF
import os

BASE_DIR = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
PRES_DIR = os.path.join(BASE_DIR, 'presentation')
os.makedirs(PRES_DIR, exist_ok=True)

def generate_pptx():
    prs = Presentation()
    
    # 1. Title Slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Mini ML Use Case: Customer Churn Prediction"
    subtitle.text = "End-to-End Business Intelligence Pipeline\nInternship Capstone Project"
    
    # 2. Business Problem
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Business Problem & Objectives"
    content = slide.placeholders[1].text_frame
    content.text = "Problem Statement: High customer attrition reduces overall revenue and LTV."
    p = content.add_paragraph()
    p.text = "Objectives:"
    p.level = 0
    content.add_paragraph().text = "1. Accurately predict which customers are at risk of churning."
    content.add_paragraph().text = "2. Identify key drivers of churn (e.g., Contracts, Service types)."
    content.add_paragraph().text = "3. Provide actionable retention recommendations."
    
    # 3. Model Results
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Machine Learning Results"
    content = slide.placeholders[1].text_frame
    content.text = "Model Evaluation:"
    content.add_paragraph().text = "- Gradient Boosting / XGBoost consistently outperformed simple linear models."
    content.add_paragraph().text = "- Achieved high ROC-AUC (~0.85+), balancing Precision and Recall effectively."
    content.add_paragraph().text = "- Pipeline successfully automated Missing Value Imputation, OHE, and Scaling."
    
    # 4. SHAP Insights
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Explainable AI (SHAP Insights)"
    content = slide.placeholders[1].text_frame
    content.text = "Key Drivers of Churn:"
    content.add_paragraph().text = "- Contract Type: Month-to-month contracts are the highest risk factor."
    content.add_paragraph().text = "- Tenure: Newer customers churn significantly more. Risk drops after 12 months."
    content.add_paragraph().text = "- Services: Fiber optic users without Tech Support exhibit high churn rates."
    
    # 5. Recommendations
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Business Recommendations & ROI"
    content = slide.placeholders[1].text_frame
    content.text = "Strategic Actions:"
    content.add_paragraph().text = "- Incentive Programs: Offer discounted first 3 months to push users to 1-year contracts."
    content.add_paragraph().text = "- Service Bundling: Bundle Tech Support with Fiber Optic packages at a discount."
    content.add_paragraph().text = "- Proactive Outreach: Target 'High Risk' flagged customers 2 weeks before billing cycles."
    
    prs.save(os.path.join(PRES_DIR, 'project_summary.pptx'))
    print("Presentation generated: project_summary.pptx")

def generate_pdf():
    pdf = FPDF()
    pdf.add_page()
    pdf.set_font("Arial", size=15, style='B')
    pdf.cell(200, 10, text="Capstone Project Summary: Customer Churn Prediction", new_x="LMARGIN", new_y="NEXT", align='C')
    
    pdf.set_font("Arial", size=12)
    pdf.ln(10)
    pdf.multi_cell(0, 10, text="1. Business Problem: Telecom companies face immense revenue loss due to customer churn. It is 5x more expensive to acquire a new customer than to retain an existing one. Our goal is to predict churn accurately.")
    
    pdf.ln(5)
    pdf.multi_cell(0, 10, text="2. Workflow: We engineered meaningful business features (Tenure Groups, Spending Categories, Risk Indicators). We handled data preprocessing strictly inside Scikit-Learn Pipelines to prevent data leakage.")
    
    pdf.ln(5)
    pdf.multi_cell(0, 10, text="3. Results: XGBoost/Random Forest emerged as top predictors. Through Hyperparameter tuning via GridSearchCV, we maximized ROC-AUC to ensure robust identification of churn candidates.")
    
    pdf.ln(5)
    pdf.multi_cell(0, 10, text="4. Explainable AI: Using SHAP values, we decoded the Blackbox model. Month-to-month contracts, short tenure, and lack of technical support are the primary catalysts for churn.")
    
    pdf.ln(5)
    pdf.multi_cell(0, 10, text="5. Recommendations: The business should immediately target the Month-to-Month segment with loyalty lock-in incentives and deploy the predict.py script in a daily batch pipeline to flag at-risk accounts.")
    
    pdf.output(os.path.join(PRES_DIR, 'project_summary.pdf'))
    print("PDF generated: project_summary.pdf")

if __name__ == "__main__":
    generate_pptx()
    generate_pdf()
